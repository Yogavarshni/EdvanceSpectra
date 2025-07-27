from google.adk.agents import Agent
from google.adk.tools import FunctionTool
from dotenv import load_dotenv
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
import os
import uuid

import firebase_admin
from firebase_admin import credentials, storage as firebase_storage
from google.cloud import storage as gcs_storage

# Load environment variables
load_dotenv()

MODEL = "gemini-2.5-pro"

# === FIREBASE CONFIGURATION ===
FIREBASE_CRED_PATH = os.path.join("orchestrationagent", "edsphere-ai-firebase-adminsdk-fbsvc-30700780ed.json")
STORAGE_BUCKET = "edsphere-ai"  # Must match Firebase bucket name # Bucket must end with .appspot.com

# === INITIALIZE FIREBASE ===
if not firebase_admin._apps:
    cred = credentials.Certificate(FIREBASE_CRED_PATH)
    firebase_admin.initialize_app(cred, {
        'storageBucket': STORAGE_BUCKET
    })
bucket = firebase_storage.bucket()

load_dotenv()

MODEL = "gemini-2.5-pro"

# ---------- TOOL FUNCTION: PDF ----------
def save_text_as_pdf(text: str) -> dict:
    try:
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        lines = text.split("\n")
        for line in lines:
            pdf.multi_cell(0, 10, line)

        file_name = f"{uuid.uuid4().hex}.pdf"
        file_path = os.path.join("generated", file_name)
        os.makedirs("generated", exist_ok=True)
        pdf.output(file_path)

        blob_name = f"materials/pdfs/{file_name}"
        blob = bucket.blob(blob_name)
        blob.upload_from_filename(file_path)
        blob.make_public()
        firebase_url = blob.public_url

        return {
            "status": "success",
            "message": "PDF Material saved and uploaded to Firebase.",
            "local_filepath": file_path,
            "firebase_url": firebase_url
        }

    except Exception as e:
        return {
            "status": "error",
            "message": f"Failed to save or upload PDF: {str(e)}"
        }

# ---------- TOOL FUNCTION: PPT ----------
def save_text_as_ppt(text: str) -> dict:
    try:
        prs = Presentation()
        bullet_slide_layout = prs.slide_layouts[1]

        slides = text.split("Slide ")
        for slide_content in slides:
            slide_content = slide_content.strip()
            if not slide_content:
                continue

            lines = slide_content.split("\n")
            title = lines[0] if lines else "Untitled"
            content = "\n".join(lines[1:]) if len(lines) > 1 else ""

            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

        file_name = f"ppt_{uuid.uuid4().hex}.pptx"
        file_path = os.path.join("generated_materials/ppts", file_name)
        os.makedirs("generated_materials/ppts", exist_ok=True)
        prs.output(file_path)

        blob_name = f"materials/ppts/{file_name}"
        blob = bucket.blob(blob_name)
        blob.upload_from_filename(file_path)
        blob.make_public()
        firebase_url = blob.public_url

        return {
            "status": "success",
            "message": "PPT Material saved and uploaded to Firebase.",
            "local_filepath": file_path,
            "firebase_url": firebase_url
        }

    except Exception as e:
        return {
            "status": "error",
            "message": f"Failed to save or upload PPT: {str(e)}"
        }

# ---------- REGISTER TOOLS ----------
save_pdf_tool = FunctionTool(func=save_text_as_pdf)
save_ppt_tool = FunctionTool(func=save_text_as_ppt)

# ---------- MATERIAL GENERATION AGENT ----------
material_generator_agent = Agent(
    name="material_generator_agent",
    model=MODEL,
    description="Generates study materials such as PDFs and PPTs from given content.",
    instruction=(
    "You are an assistant that generates structured study material from a topic and grade level.\n\n"
    "ALWAYS follow this format strictly. No deviations allowed.\n\n"
    "Output must begin with:\n"
    "Title: <Title of the topic> (Grade <grade number>)\n\n"
    "Then follow with slide content:\n"
    "Slide 1: <Slide Title>\n"
    "- Bullet point 1\n"
    "- Bullet point 2\n\n"
    "Slide 2: <Slide Title>\n"
    "- Bullet point 1\n"
    "- Bullet point 2\n"
    "...\n\n"
    "Once all slides are generated, END by calling either:\n"
    "`save_pdf_tool` OR `save_ppt_tool`, passing the entire material string as input.\n\n"
    "VERY IMPORTANT:\n"
    "- Do NOT include any extra explanation.\n"
    "- Do NOT continue writing after the tool call.\n"
    "- NEVER skip the tool call.\n"
    "- You MUST treat the tool call as the final step.\n"
    "- ONLY use `save_pdf_tool` or `save_ppt_tool`, not both.\n"
    "Return both the local file path and Firebase public URL, then stop."
    ),
    tools=[save_pdf_tool, save_ppt_tool]
)
